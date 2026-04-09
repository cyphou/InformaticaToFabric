#!/usr/bin/env python3
"""Generate a PowerPoint presentation documenting the Informatica-to-Fabric migration tool and its ROI."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── Brand colours ─────────────────────────────────────────────
FABRIC_BLUE   = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE     = RGBColor(0x00, 0x3C, 0x6E)
ACCENT_ORANGE = RGBColor(0xFF, 0x45, 0x00)
ACCENT_GREEN  = RGBColor(0x27, 0xAE, 0x60)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
DARK_GRAY     = RGBColor(0x2C, 0x3E, 0x50)
LIGHT_GRAY    = RGBColor(0xEC, 0xF0, 0xF1)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
BLACK         = RGBColor(0x00, 0x00, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _add_bg(slide, color):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER):
    """Add a rounded-rectangle shape with centred text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = align
        run = tf.paragraphs[0].add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.bold = bold
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18, font_color=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16, font_color=DARK_GRAY, spacing=Pt(6)):
    """Add a bullet-point list inside a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        p.level = 0
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
    return txBox


def _add_title_bar(slide, title, subtitle=None):
    """Dark-blue title bar across top of slide."""
    _add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.3), DARK_BLUE, bold=True)
    _add_text_box(slide, Inches(0.6), Inches(0.2), Inches(11), Inches(0.6), title, font_size=28, font_color=WHITE, bold=True)
    if subtitle:
        _add_text_box(slide, Inches(0.6), Inches(0.75), Inches(11), Inches(0.5), subtitle, font_size=16, font_color=RGBColor(0xBD, 0xC3, 0xC7))


def _add_kpi_card(slide, left, top, value, label, color):
    """Small KPI card (value + label)."""
    _add_shape(slide, left, top, Inches(2.2), Inches(1.5), color)
    _add_text_box(slide, left + Inches(0.05), top + Inches(0.15), Inches(2.1), Inches(0.8), value, font_size=32, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _add_text_box(slide, left + Inches(0.05), top + Inches(0.85), Inches(2.1), Inches(0.5), label, font_size=12, font_color=WHITE, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════
#  SLIDES
# ═════════════════════════════════════════════════════════════════

def slide_title(prs):
    """Slide 1 — Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _add_bg(slide, DARK_BLUE)

    _add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
                  "Informatica PowerCenter / IDMC  \u2192  Microsoft Fabric / Azure Databricks",
                  font_size=36, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(0.8),
                  "Automated Migration Tool \u2014 PowerCenter + IICS + IDMC (12 Services) \u2192 Fabric / Databricks / DBT",
                  font_size=22, font_color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)

    # Stat bar
    stats = [
        ("97 / 100", "Sprints Complete", ACCENT_ORANGE),
        ("1,843", "Automated Tests", ACCENT_GREEN),
        ("6", "AI Agents", FABRIC_BLUE),
        ("3", "Target Platforms", ACCENT_PURPLE),
    ]
    x = Inches(1.5)
    for val, lbl, clr in stats:
        _add_kpi_card(slide, x, Inches(4.0), val, lbl, clr)
        x += Inches(2.6)

    _add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
                  "Powered by GitHub Copilot  •  6-Agent AI Architecture  •  Zero Manual Code Rewriting",
                  font_size=14, font_color=RGBColor(0x95, 0xA5, 0xA6), align=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 2 — The Problem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "The Challenge", "Why migrate away from Informatica?")

    pain_points = [
        "Informatica PowerCenter/IICS/IDMC licensing costs escalate year-over-year (avg. 15-25% annual increase)",
        "On-premises infrastructure limits elasticity \u2014 peak workloads pay for idle capacity",
        "Vendor lock-in: proprietary XML formats, no native cloud integration, limited Spark support",
        "IDMC cloud services (CDI, CDGC, CDQ, MDM, etc.) add complexity without reducing on-prem footprint",
        "Manual migration projects historically take 12-24 months with 40-60% cost overruns",
        "Lack of modern data engineering patterns (Medallion, Delta Lake, Unity Catalog)",
        "AutoSys / scheduler coupling adds hidden dependency & operational risk",
        "Skills scarcity: Informatica specialists are increasingly hard to hire",
    ]
    _add_bullet_list(slide, Inches(0.6), Inches(1.6), Inches(7.5), Inches(5.5), pain_points, font_size=16)

    # Side callout
    _add_shape(slide, Inches(8.8), Inches(1.8), Inches(4), Inches(2.2), ACCENT_ORANGE,
               "Industry Fact:\n\nEnterprise Informatica\nlicense spend averages\n$1.2–3.5M / year",
               font_size=15, font_color=WHITE, bold=False)

    _add_shape(slide, Inches(8.8), Inches(4.3), Inches(4), Inches(2.2), RGBColor(0xC0, 0x39, 0x2B),
               "Manual Migration:\n\n12–24 months\n40–60% cost overrun\n6–12 FTEs required",
               font_size=15, font_color=WHITE, bold=False)


def slide_solution(prs):
    """Slide 3 — Solution Overview."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "The Solution", "AI-Orchestrated, Fully Automated Migration")

    _add_text_box(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.6),
                  "A 6-agent Copilot architecture that converts Informatica XML exports end-to-end, "
                  "producing production-ready Fabric/Databricks artifacts with zero manual code rewriting.",
                  font_size=16, font_color=DARK_GRAY)

    # Pipeline flow
    labels = [
        ("📂 XML\nExports", ACCENT_ORANGE),
        ("🔍 Assess\n& Inventory", RGBColor(0xE6, 0x7E, 0x22)),
        ("🗄️ SQL\nConvert", ACCENT_PURPLE),
        ("📓 Notebook\nGenerate", ACCENT_GREEN),
        ("⚡ Pipeline\nGenerate", FABRIC_BLUE),
        ("✅ Validate\n& Deploy", RGBColor(0xC0, 0x39, 0x2B)),
    ]
    x = Inches(0.4)
    for i, (lbl, clr) in enumerate(labels):
        _add_shape(slide, x, Inches(2.4), Inches(1.9), Inches(1.1), clr, lbl, font_size=13, font_color=WHITE, bold=True)
        if i < len(labels) - 1:
            # Proper arrow shape instead of text
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, x + Inches(1.93), Inches(2.65), Inches(0.28), Inches(0.35))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()
        x += Inches(2.15)

    # Key features below
    features_l = [
        "PowerCenter 9.x/10.x + IICS + IDMC (12 cloud services) full support",
        "6 source databases (Oracle, SQL Server, Teradata, DB2, MySQL, PostgreSQL)",
        "180+ SQL pattern conversions → Spark SQL",
        "18 transformation types → PySpark code",
        "AutoSys JIL → Pipeline / Workflow JSON",
        "Medallion architecture (Bronze / Silver / Gold)",
    ]
    features_r = [
        "Delta Lake DDL with partition strategy",
        "Plugin system for custom enterprise rules",
        "Python SDK + REST API for CI/CD",
        "11-level validation framework + statistical checks",
        "Datadog observability (logs, metrics, APM tracing)",
        "Agentic alerting with auto-remediation & learning",
        "IDMC full component assessment (12 services)",
        "Data catalog integration (Purview / Unity Catalog)",
        "IaC generation (Terraform + Bicep)",
        "Container & K8s deployment (Helm chart)",
        "CI/CD pipelines (GitHub Actions + Azure DevOps)",
    ]
    _add_bullet_list(slide, Inches(0.4), Inches(3.9), Inches(6), Inches(3.3), features_l, font_size=14)
    _add_bullet_list(slide, Inches(6.7), Inches(3.9), Inches(6), Inches(3.3), features_r, font_size=14)


def slide_architecture(prs):
    """Slide 4 — Architecture."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Multi-Agent Architecture", "6 specialized AI agents orchestrated by GitHub Copilot")

    agents = [
        ("🎯 Orchestrator", "Plans & coordinates\nmigration waves", DARK_BLUE),
        ("🔍 Assessment", "XML parsing, inventory,\ncomplexity scoring", RGBColor(0xE6, 0x7E, 0x22)),
        ("🗄️ SQL Migration", "Oracle/MSSQL/Teradata\n→ Spark SQL / T-SQL", ACCENT_PURPLE),
        ("📓 Notebook Migration", "Mappings → PySpark\nnotebooks / DBT", ACCENT_GREEN),
        ("⚡ Pipeline Migration", "Workflows → Fabric\nPipeline / DB Workflow", FABRIC_BLUE),
        ("✅ Validation", "Test generation,\nrow counts & checksums", RGBColor(0xC0, 0x39, 0x2B)),
    ]
    x = Inches(0.3)
    for name, desc, clr in agents:
        _add_shape(slide, x, Inches(1.6), Inches(2.0), Inches(1.0), clr, name, font_size=14, font_color=WHITE, bold=True)
        _add_text_box(slide, x + Inches(0.05), Inches(2.7), Inches(1.9), Inches(0.9), desc, font_size=11, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)
        x += Inches(2.15)

    # Targets
    _add_text_box(slide, Inches(0.6), Inches(3.9), Inches(12), Inches(0.4),
                  "Target Platforms", font_size=20, font_color=DARK_BLUE, bold=True)

    targets = [
        ("Microsoft Fabric", "Notebooks (notebookutils)\nData Pipelines JSON\nLakehouse DDL\nOneLake + Medallion", FABRIC_BLUE),
        ("Azure Databricks", "Notebooks (dbutils)\nWorkflows (Jobs API)\nUnity Catalog DDL\nDLT + Delta Sharing", RGBColor(0xFF, 0x36, 0x21)),
        ("DBT on Databricks", "staging / intermediate / marts\nJinja SQL models\nCI/CD pipeline\nSources + Schema YAML", ACCENT_PURPLE),
    ]
    x = Inches(0.5)
    for name, desc, clr in targets:
        _add_shape(slide, x, Inches(4.5), Inches(3.8), Inches(0.7), clr, name, font_size=16, font_color=WHITE, bold=True)
        _add_text_box(slide, x + Inches(0.1), Inches(5.25), Inches(3.6), Inches(1.8), desc, font_size=12, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)
        x += Inches(4.15)


def slide_what_gets_migrated(prs):
    """Slide 5 — What Gets Migrated."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "What Gets Migrated", "Comprehensive Informatica coverage")

    cols = [
        ("Mappings → Notebooks", [
            "Source Qualifier, Expression, Filter",
            "Aggregator, Joiner, Lookup, Router",
            "Update Strategy, Rank, Union",
            "Normalizer, Sorter, Sequence Gen",
            "Stored Procedure, Mapplets",
            "SQL Transformation, Data Masking",
        ], ACCENT_GREEN),
        ("Workflows → Pipelines", [
            "Sessions, Command Tasks, Timers",
            "Decisions, Event Wait/Raise",
            "Assignments, Email Tasks, Worklets",
            "Link Conditions, Control Tasks",
            "IICS Taskflows (cloud-native)",
            "Schedule triggers → cron expressions",
        ], FABRIC_BLUE),
        ("SQL → Spark SQL", [
            "Oracle: NVL, DECODE, SYSDATE, ROWNUM",
            "CONNECT BY, (+) joins, analytics",
            "SQL Server: GETDATE, CROSS APPLY",
            "Teradata: QUALIFY, SAMPLE, ZEROIFNULL",
            "DB2: FETCH FIRST, VALUE, RRN",
            "MySQL/PostgreSQL: casts, arrays",
        ], ACCENT_PURPLE),
        ("AutoSys JIL", [
            "BOX → Pipeline container",
            "CMD (pmcmd) → Notebook activity",
            "FW → File sensor / trigger",
            "Conditions (s/f/n/d) → dependsOn",
            "Calendars → cron schedules",
            "Cross-box → Execute Pipeline refs",
        ], ACCENT_ORANGE),
    ]

    x = Inches(0.3)
    for title, items, clr in cols:
        _add_shape(slide, x, Inches(1.5), Inches(3.05), Inches(0.6), clr, title, font_size=14, font_color=WHITE, bold=True)
        _add_bullet_list(slide, x + Inches(0.1), Inches(2.15), Inches(2.85), Inches(4.8), items, font_size=12, font_color=DARK_GRAY, spacing=Pt(4))
        x += Inches(3.2)


def slide_sql_depth(prs):
    """Slide 6 — Advanced SQL & PL/SQL Capabilities."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Advanced SQL & PL/SQL Engine", "Beyond simple find-and-replace — structured conversion")

    items_l = [
        "180+ regex conversion patterns across 6 database dialects",
        "Cursor → PySpark iterator (explicit, implicit, FOR...IN)",
        "BULK COLLECT INTO → DataFrame spark.sql()",
        "FORALL INSERT/UPDATE → batch DataFrame write / Delta MERGE",
        "Exception blocks → try/except with mapped Python exceptions",
        "Package state variables → Python module-level state",
        "EXECUTE IMMEDIATE → spark.sql() with parameter substitution",
        "CONNECT BY PRIOR → WITH RECURSIVE CTE",
    ]
    items_r = [
        "PIVOT retains native Spark SQL / UNPIVOT → stack()",
        "Correlated EXISTS → LEFT SEMI JOIN",
        "Correlated NOT EXISTS → LEFT ANTI JOIN",
        "Temporal tables (FOR SYSTEM_TIME) → Delta time-travel",
        "GTT, Materialized Views, DB Link detection & flagging",
        "DECODE → CASE WHEN expansion",
        "Performance annotations (-- PERF: comments on anti-patterns)",
        "Configurable YAML/JSON rule engine for enterprise overrides",
    ]
    _add_bullet_list(slide, Inches(0.4), Inches(1.6), Inches(6), Inches(5.5), items_l, font_size=14)
    _add_bullet_list(slide, Inches(6.7), Inches(1.6), Inches(6), Inches(5.5), items_r, font_size=14)


def slide_validation(prs):
    """Slide 7 — Validation & Quality."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Validation & Quality Assurance", "11-level automated validation framework")

    levels = [
        ("L1", "Row Count", "Source vs Target row count match"),
        ("L2", "Column Checksum", "Hash-based column integrity check"),
        ("L3", "Aggregate Check", "SUM / AVG / MIN / MAX comparison"),
        ("L4", "Sample Diff", "Row-by-row sample record comparison"),
        ("L5", "NULL / Duplicate", "NULL distribution & duplicate detection"),
        ("L6", "Statistical", "Mean/stddev distribution + K-S test"),
        ("L7", "SCD Type 2", "Effective dates, current_flag, gap/overlap check"),
        ("L8", "Null Distribution", "Column-level null % comparison (1% threshold)"),
        ("L9", "Referential Integrity", "FK relationship validation from LKP/JNR"),
        ("L10", "A/B Testing", "Side-by-side source/target row comparison"),
        ("L11", "Business Rules", "Custom assertions from migration.yaml"),
    ]

    y = Inches(1.5)
    for lvl, name, desc in levels:
        clr = ACCENT_GREEN if int(lvl[1:]) <= 5 else FABRIC_BLUE if int(lvl[1:]) <= 8 else ACCENT_PURPLE
        _add_shape(slide, Inches(0.4), y, Inches(0.7), Inches(0.42), clr, lvl, font_size=11, font_color=WHITE, bold=True)
        _add_text_box(slide, Inches(1.2), y + Inches(0.02), Inches(2.2), Inches(0.4), name, font_size=13, font_color=DARK_BLUE, bold=True)
        _add_text_box(slide, Inches(3.5), y + Inches(0.02), Inches(8), Inches(0.4), desc, font_size=13, font_color=DARK_GRAY)
        y += Inches(0.48)

    _add_text_box(slide, Inches(0.4), Inches(6.9), Inches(12), Inches(0.4),
                  "1,843 automated tests  •  Purview & Unity Catalog lineage  •  Column-level impact analysis",
                  font_size=14, font_color=RGBColor(0x95, 0xA5, 0xA6), align=PP_ALIGN.CENTER)


def slide_extensibility(prs):
    """Slide 8 — Extensibility & SDK."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Extensibility & Enterprise Integration", "Plugin system, SDK, REST API, rule engine")

    cards = [
        ("Plugin System", [
            "@register_transform decorator",
            "@register_sql_rewrite decorator",
            "@post_notebook / @post_pipeline hooks",
            "Auto-discovery from plugins/ directory",
            "3 example plugins included",
        ], ACCENT_PURPLE),
        ("Python SDK", [
            "MigrationConfig dataclass",
            "MigrationSDK.migrate() / .assess()",
            "Convenience: from sdk import migrate",
            "YAML / dict configuration",
            "Programmatic CI/CD integration",
        ], ACCENT_GREEN),
        ("REST API", [
            "POST /migrate → async job",
            "POST /assess → inventory",
            "GET /status/{job_id}",
            "FastAPI (or stdlib fallback)",
            "OpenAPI / Swagger auto-docs",
        ], FABRIC_BLUE),
        ("Rule Engine", [
            "YAML / JSON rule definitions",
            "Priority-based application",
            "Enterprise override merging",
            "Default rules/ directory",
            "10+ SQL rules out-of-box",
        ], ACCENT_ORANGE),
    ]

    x = Inches(0.3)
    for title, items, clr in cards:
        _add_shape(slide, x, Inches(1.5), Inches(3.05), Inches(0.55), clr, title, font_size=15, font_color=WHITE, bold=True)
        _add_bullet_list(slide, x + Inches(0.1), Inches(2.1), Inches(2.85), Inches(4.5), items, font_size=12, font_color=DARK_GRAY, spacing=Pt(4))
        x += Inches(3.2)


def slide_observability(prs):
    """Slide — Observability & Intelligent Monitoring."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Observability & Intelligent Monitoring",
                   "Datadog integration, agentic alerting, global monitoring platform")

    cards = [
        ("Datadog Integration", [
            "Structured logs → Datadog Log Explorer",
            "Custom metrics per phase & artifact type",
            "APM distributed tracing (ddtrace)",
            "Monitor definitions for SLA breaches",
            "Graceful degradation if not installed",
        ], RGBColor(0x63, 0x2C, 0xA6)),
        ("Agentic Alerting", [
            "Signal processor with confidence scoring",
            "Auto-remediation actions (3 modes)",
            "Learning loop with SQLite history",
            "Pattern recognition → auto-fix",
            "Escalation when confidence < threshold",
        ], ACCENT_ORANGE),
        ("Global Monitoring", [
            "Unified control plane for all targets",
            "4-tier escalation chains (L1→L4)",
            "SLO tracking & compliance reporting",
            "Enterprise dashboards (Ops + Exec)",
            "Cross-component alert correlation",
        ], FABRIC_BLUE),
        ("IDMC Assessment", [
            "12 IDMC components parsed via REST",
            "CDI / CDGC / CDQ / MDM / DI / B2B",
            "Migration review: merge/optimize/rework",
            "Quality scoring & readiness gate",
            "IDMC-to-Fabric comparison report",
        ], ACCENT_GREEN),
    ]

    x = Inches(0.3)
    for title, items, clr in cards:
        _add_shape(slide, x, Inches(1.5), Inches(3.05), Inches(0.55), clr, title, font_size=15, font_color=WHITE, bold=True)
        _add_bullet_list(slide, x + Inches(0.1), Inches(2.1), Inches(2.85), Inches(4.5), items, font_size=12, font_color=DARK_GRAY, spacing=Pt(4))
        x += Inches(3.2)


def slide_roi_model(prs):
    """Slide 9 — ROI Model."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "ROI Analysis", "Return on Investment — Automated vs Manual Migration")

    # --- Left column: cost model table ---
    _add_text_box(slide, Inches(0.5), Inches(1.5), Inches(6), Inches(0.4),
                  "Cost Comparison  (100-mapping migration, typical enterprise)", font_size=16, font_color=DARK_BLUE, bold=True)

    rows = [
        ("", "Manual Migration", "Automated (This Tool)"),
        ("Duration", "12–18 months", "4–8 weeks"),
        ("FTEs Required", "6–12 engineers", "1–2 engineers"),
        ("Avg FTE Cost (annual)", "$150K–$200K", "$150K–$200K"),
        ("Total Labor Cost", "$900K – $2.4M", "$58K – $154K"),
        ("Informatica License (saved)", "—", "$1.2M – $3.5M / year"),
        ("Rework / Bug-Fix", "20–30% of budget", "< 5% (test-validated)"),
        ("Validation Effort", "Manual QA: 2–4 months", "Auto-generated: 1–2 days"),
    ]

    y = Inches(2.0)
    for i, (label, manual, auto) in enumerate(rows):
        bg = DARK_BLUE if i == 0 else LIGHT_GRAY if i % 2 == 0 else WHITE
        fc = WHITE if i == 0 else DARK_GRAY
        bld = i == 0
        _add_shape(slide, Inches(0.5), y, Inches(2.5), Inches(0.45), bg, label, font_size=12, font_color=fc, bold=bld, align=PP_ALIGN.LEFT)
        _add_shape(slide, Inches(3.0), y, Inches(2.5), Inches(0.45), bg, manual, font_size=12, font_color=fc, bold=bld, align=PP_ALIGN.CENTER)
        _add_shape(slide, Inches(5.5), y, Inches(2.7), Inches(0.45), bg if i == 0 else ACCENT_GREEN if i > 0 else bg,
                   auto, font_size=12, font_color=WHITE if i > 0 else fc, bold=bld, align=PP_ALIGN.CENTER)
        y += Inches(0.48)

    # --- Right column: KPI cards ---
    _add_kpi_card(slide, Inches(8.7), Inches(1.6), "85–95%", "Cost Reduction\nvs Manual Migration", ACCENT_GREEN)
    _add_kpi_card(slide, Inches(8.7), Inches(3.4), "3–6×", "Faster Time-\nto-Production", FABRIC_BLUE)
    _add_kpi_card(slide, Inches(8.7), Inches(5.2), "$1.2–3.5M", "Annual License\nSavings", ACCENT_ORANGE)

    _add_kpi_card(slide, Inches(11.0), Inches(1.6), "80–90%", "Reduction in\nFTE Requirement", ACCENT_PURPLE)
    _add_kpi_card(slide, Inches(11.0), Inches(3.4), "99%+", "Automated Test\nCoverage", RGBColor(0xC0, 0x39, 0x2B))
    _add_kpi_card(slide, Inches(11.0), Inches(5.2), "< 5%", "Rework Rate\nvs 20-30% Manual", DARK_BLUE)


def slide_roi_timeline(prs):
    """Slide 10 — ROI Timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Time-to-Value", "Manual vs Automated Migration Timeline")

    # Manual timeline bar
    _add_text_box(slide, Inches(0.5), Inches(1.6), Inches(4), Inches(0.4),
                  "Manual Migration (Traditional)", font_size=16, font_color=DARK_GRAY, bold=True)

    manual_phases = [
        ("Discovery\n& Planning", "3 mo", RGBColor(0x95, 0xA5, 0xA6)),
        ("SQL\nConversion", "4 mo", RGBColor(0x95, 0xA5, 0xA6)),
        ("Notebook\nDevelopment", "4 mo", RGBColor(0x95, 0xA5, 0xA6)),
        ("Testing\n& QA", "3 mo", RGBColor(0x95, 0xA5, 0xA6)),
        ("Deploy &\nStabilize", "4 mo", RGBColor(0x95, 0xA5, 0xA6)),
    ]
    x = Inches(0.5)
    for name, dur, clr in manual_phases:
        _add_shape(slide, x, Inches(2.1), Inches(2.35), Inches(0.85), clr, f"{name}\n{dur}", font_size=11, font_color=WHITE, bold=False)
        x += Inches(2.45)

    _add_shape(slide, Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.35), RGBColor(0xC0, 0x39, 0x2B),
               "Total: 12–18 months  •  6–12 FTEs  •  $900K–$2.4M labor", font_size=12, font_color=WHITE, bold=True)

    # Automated timeline bar
    _add_text_box(slide, Inches(0.5), Inches(3.8), Inches(4), Inches(0.4),
                  "Automated Migration (This Tool)", font_size=16, font_color=DARK_GRAY, bold=True)

    auto_phases = [
        ("Assessment\n(automated)", "1 day", ACCENT_GREEN),
        ("Conversion\n(automated)", "1–2 days", ACCENT_GREEN),
        ("Review &\nCustomize", "2–4 weeks", FABRIC_BLUE),
        ("Validate\n& Deploy", "1–2 weeks", ACCENT_GREEN),
    ]
    x = Inches(0.5)
    for name, dur, clr in auto_phases:
        w = Inches(1.6) if "day" in dur and "week" not in dur else Inches(3.6) if "4 week" in dur else Inches(2.8)
        _add_shape(slide, x, Inches(4.3), w, Inches(0.85), clr, f"{name}\n{dur}", font_size=11, font_color=WHITE, bold=False)
        x += w + Inches(0.15)

    _add_shape(slide, Inches(0.5), Inches(5.25), Inches(8.8), Inches(0.35), ACCENT_GREEN,
               "Total: 4–8 weeks  •  1–2 FTEs  •  $58K–$154K labor", font_size=12, font_color=WHITE, bold=True)

    # Savings callout
    _add_shape(slide, Inches(9.6), Inches(4.3), Inches(3.4), Inches(2.5), DARK_BLUE,
               "Net Savings\n\n$840K – $2.25M\nlabor cost avoided\n\n+ $1.2–3.5M/year\nlicense elimination\n\nPayback: < 3 months",
               font_size=15, font_color=WHITE, bold=False)


def slide_customer_scenarios(prs):
    """Slide 11 — Customer Scenarios."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Enterprise Scenarios & Scale", "Real-world migration complexity handled")

    scenarios = [
        ("Mid-Size Enterprise", "50–100 mappings\n20–30 workflows\n1–2 database types\n\nTimeline: 3–4 weeks\nFTEs: 1 engineer\nSavings: $400K–$800K", FABRIC_BLUE),
        ("Large Enterprise", "200–500 mappings\n50–100 workflows\n3–4 database types\nAutoSys JIL jobs\n\nTimeline: 6–10 weeks\nFTEs: 2 engineers\nSavings: $1.5M–$3M", ACCENT_ORANGE),
        ("Global Financial Org", "1,000+ mappings\n200+ workflows\n6 database types\nAutoSys + PL/SQL heavy\n\nTimeline: 12–16 weeks\nFTEs: 3–4 engineers\nSavings: $4M–$8M", ACCENT_PURPLE),
    ]

    x = Inches(0.5)
    for title, desc, clr in scenarios:
        _add_shape(slide, x, Inches(1.5), Inches(3.8), Inches(0.6), clr, title, font_size=16, font_color=WHITE, bold=True)
        _add_text_box(slide, x + Inches(0.15), Inches(2.2), Inches(3.5), Inches(4.5), desc, font_size=13, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)
        x += Inches(4.15)


def slide_security(prs):
    """Slide 12 — Security & Governance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Security, Governance & Compliance", "Enterprise-grade safeguards built-in")

    items = [
        "Credential sanitization — passwords/keys stripped from all generated artifacts",
        "PII detection — flags columns containing personal data (names, SSN, email patterns)",
        "Audit logging — every migration action recorded in JSON audit trail",
        "Multi-tenant Key Vault integration — secrets resolved via notebookutils / dbutils",
        "Data Quality rules — DQ assertions auto-generated from mapping metadata",
        "Column-level lineage — full source → transform → target traceability",
        "Microsoft Purview integration — Apache Atlas v2 metadata entities + lineage",
        "Unity Catalog tags — ALTER TABLE SET TAGS for migration provenance",
        "Impact analysis — trace downstream effects of source table changes",
        "5-level validation framework — statistical verification of migrated data",
    ]
    _add_bullet_list(slide, Inches(0.5), Inches(1.6), Inches(7.5), Inches(5.5), items, font_size=15)

    _add_shape(slide, Inches(8.8), Inches(1.8), Inches(4), Inches(4.5), DARK_BLUE,
               "Compliance Coverage\n\n✓ SOX audit trail\n✓ GDPR PII detection\n✓ HIPAA data masking flags\n✓ Zero credential exposure\n✓ Role-based access via\n   Key Vault / Unity Catalog\n✓ Full lineage for regulators",
               font_size=14, font_color=WHITE, bold=False)


def slide_cloud_native(prs):
    """Slide — Cloud-Native DevOps & Scale."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Cloud-Native DevOps & Scale",
                   "Infrastructure as Code, containerized deployment, CI/CD pipelines, and enterprise benchmarks")

    # Three columns: IaC, Container/K8s, CI/CD
    cols = [
        ("🏗️ Infrastructure as Code", ACCENT_PURPLE, [
            "Terraform HCL for Fabric & Databricks",
            "Azure Bicep templates with parameters",
            "Multi-environment variable extraction",
            "Validated output (HCL lint + Bicep lint)",
        ]),
        ("🐳 Container & Kubernetes", FABRIC_BLUE, [
            "Production Dockerfile (Python 3.12-slim)",
            "Docker Compose (API + Web + Redis)",
            "K8s manifests (Deployment, Service, ConfigMap)",
            "Helm chart with configurable values",
        ]),
        ("🔄 CI/CD Pipelines", ACCENT_GREEN, [
            "GitHub Actions workflow generation",
            "Azure DevOps multi-stage pipelines",
            "Environment gates (dev → test → prod)",
            "Secret injection (never hardcoded)",
        ]),
    ]

    x = Inches(0.3)
    for title, clr, items in cols:
        _add_shape(slide, x, Inches(1.6), Inches(4.0), Inches(0.6), clr, title,
                   font_size=15, font_color=WHITE, bold=True)
        _add_bullet_list(slide, x + Inches(0.1), Inches(2.3), Inches(3.8), Inches(2.5),
                         items, font_size=13)
        x += Inches(4.3)

    # Benchmark row
    _add_shape(slide, Inches(0.3), Inches(5.0), Inches(12.5), Inches(0.5), ACCENT_ORANGE,
               "📊 Enterprise Benchmarks: 500+ mapping stress tests  •  Memory profiling  •  "
               "Parallel generation (ProcessPoolExecutor)  •  Golden dataset regression suite",
               font_size=13, font_color=WHITE, bold=False)

    _add_text_box(slide, Inches(0.4), Inches(5.7), Inches(12), Inches(0.8),
                  "One command deploys the entire migration tool to AKS, ECS, or any K8s cluster.\n"
                  "CI/CD pipelines automate: assess → convert → validate → deploy — with approval gates per environment.",
                  font_size=14, font_color=DARK_GRAY, align=PP_ALIGN.CENTER)


def slide_next_steps(prs):
    """Slide 13 — Next Steps / CTA."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK_BLUE)

    _add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8),
                  "Next Steps", font_size=36, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    steps = [
        ("1", "Discovery Workshop", "Inventory your Informatica estate\n(PowerCenter + IICS + AutoSys)", ACCENT_ORANGE),
        ("2", "Proof of Concept", "Migrate 5–10 representative\nmappings in 1 week", ACCENT_GREEN),
        ("3", "Wave Planning", "Topological dependency sort\n+ complexity-based waves", FABRIC_BLUE),
        ("4", "Production Migration", "Automated conversion with\nvalidation at every step", ACCENT_PURPLE),
        ("5", "Decommission Informatica", "License savings begin —\nROI realized within 3 months", RGBColor(0xC0, 0x39, 0x2B)),
    ]

    x = Inches(0.5)
    for num, title, desc, clr in steps:
        _add_shape(slide, x, Inches(2.0), Inches(2.3), Inches(1.0), clr, f"{title}", font_size=14, font_color=WHITE, bold=True)
        _add_shape(slide, x + Inches(0.7), Inches(1.55), Inches(0.6), Inches(0.5), WHITE, num, font_size=20, font_color=clr, bold=True)
        _add_text_box(slide, x + Inches(0.05), Inches(3.1), Inches(2.2), Inches(1.5), desc, font_size=12, font_color=RGBColor(0xBD, 0xC3, 0xC7), align=PP_ALIGN.CENTER)
        x += Inches(2.5)

    _add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(1.0),
                  "\"The fastest path from Informatica to Fabric/Databricks —\n"
                  "automated, validated, and production-ready.\"",
                  font_size=22, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                  "informatica-to-fabric  •  MIT Licensed  •  github.com/cyphou/InformaticaToFabric",
                  font_size=14, font_color=RGBColor(0x95, 0xA5, 0xA6), align=PP_ALIGN.CENTER)

def slide_detailed_savings(prs):
    """Slide 11 — Detailed per-task time savings breakdown."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _add_title_bar(slide, "Detailed Time & Cost Savings Estimate",
                   "Per-task breakdown — 200 mappings, 50 workflows, 3 DB types, AutoSys")

    # Table header
    headers = ["Migration Task", "Manual\n(h/item)", "Auto\n(h/item)", "Items", "Manual\nTotal (h)", "Auto\nTotal (h)", "Saved"]
    col_widths = [Inches(2.4), Inches(0.9), Inches(0.9), Inches(0.6), Inches(1.05), Inches(1.05), Inches(0.85)]

    # Data rows: (task, manual_h, auto_h, count)
    tasks = [
        ("Assessment & Inventory",             8,   0.1,  200),
        ("SQL Override Conversion",           12,   0.1,  120),
        ("Stored Procedure Conversion",       24,   0.5,   30),
        ("Mapping → PySpark Notebook",        16,   0.2,  200),
        ("Mapping → DBT Model",              10,   0.1,   80),
        ("Workflow → Pipeline JSON",          12,   0.2,   50),
        ("AutoSys JIL → Pipeline",            8,   0.1,   40),
        ("Delta Lake Schema (DDL)",            4,   0.05, 200),
        ("Validation Notebook Generation",     6,   0.1,  200),
        ("Unit/Integration Testing",          10,   0.5,  200),
        ("Lineage Documentation",              4,   0.05, 200),
        ("Data Catalog (Purview/UC) Setup",    6,   0.1,  200),
    ]

    x_off = Inches(0.3)
    y = Inches(1.55)

    # Draw header row
    x = x_off
    for hdr, w in zip(headers, col_widths):
        _add_shape(slide, x, y, w, Inches(0.55), DARK_BLUE, hdr, font_size=10, font_color=WHITE, bold=True)
        x += w

    # Draw data rows
    total_manual = 0
    total_auto = 0
    for i, (task, mh, ah, count) in enumerate(tasks):
        y += Inches(0.42)
        m_total = mh * count
        a_total = ah * count
        savings_pct = f"{((m_total - a_total) / m_total * 100):.0f}%"
        total_manual += m_total
        total_auto += a_total

        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        vals = [task, f"{mh}", f"{ah}", f"{count}", f"{m_total:,.0f}", f"{a_total:,.0f}", savings_pct]
        x = x_off
        for j, (val, w) in enumerate(zip(vals, col_widths)):
            clr = bg
            fc = DARK_GRAY
            if j == 6:  # savings column
                clr = ACCENT_GREEN
                fc = WHITE
            al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            _add_shape(slide, x, y, w, Inches(0.38), clr, val, font_size=10, font_color=fc, bold=False, align=al)
            x += w

    # Totals row
    y += Inches(0.44)
    savings_total_pct = f"{((total_manual - total_auto) / total_manual * 100):.0f}%"
    totals = ["TOTAL", "", "", "", f"{total_manual:,.0f}", f"{total_auto:,.0f}", savings_total_pct]
    x = x_off
    for j, (val, w) in enumerate(zip(totals, col_widths)):
        clr = DARK_BLUE if j < 4 else ACCENT_GREEN
        fc = WHITE
        al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        _add_shape(slide, x, y, w, Inches(0.45), clr, val, font_size=11, font_color=fc, bold=True, align=al)
        x += w

    # Summary callout on the right
    manual_fte_months = total_manual / 160  # 160h/month per FTE
    auto_fte_months = total_auto / 160
    labor_manual = total_manual * 95  # $95/h blended rate
    labor_auto = total_auto * 95
    saved = labor_manual - labor_auto

    summary = (
        f"Manual: {total_manual:,.0f} hours\n"
        f"= {manual_fte_months:.0f} FTE-months\n"
        f"= ${labor_manual/1000:,.0f}K labor\n\n"
        f"Automated: {total_auto:,.0f} hours\n"
        f"= {auto_fte_months:.1f} FTE-months\n"
        f"= ${labor_auto/1000:,.0f}K labor\n\n"
        f"Net Saved: ${saved/1000:,.0f}K\n"
        f"({savings_total_pct} reduction)"
    )
    _add_shape(slide, Inches(8.2), Inches(1.55), Inches(4.8), Inches(3.8), FABRIC_BLUE,
               summary, font_size=14, font_color=WHITE, bold=False, align=PP_ALIGN.LEFT)

    # License savings note at bottom
    _add_shape(slide, Inches(8.2), Inches(5.5), Inches(4.8), Inches(1.3), ACCENT_ORANGE,
               "Additional Savings:\n+$1.2–3.5M/year\nInformatica license elimination",
               font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════
#  MAIN
# ═════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)           # 1 - Title
    slide_problem(prs)         # 2 - The Challenge
    slide_solution(prs)        # 3 - Solution Overview
    slide_architecture(prs)    # 4 - Architecture
    slide_what_gets_migrated(prs)  # 5 - What Gets Migrated
    slide_sql_depth(prs)       # 6 - Advanced SQL
    slide_validation(prs)      # 7 - Validation
    slide_extensibility(prs)   # 8 - Extensibility
    slide_observability(prs)   # 9 - Observability & Monitoring
    slide_roi_model(prs)       # 10 - ROI Model
    slide_roi_timeline(prs)    # 11 - Timeline
    slide_detailed_savings(prs)   # 12 - Detailed Savings Breakdown
    slide_customer_scenarios(prs)  # 13 - Scenarios
    slide_security(prs)        # 14 - Security
    slide_cloud_native(prs)    # 15 - Cloud-Native DevOps & Scale
    slide_next_steps(prs)      # 16 - Next Steps

    out_path = os.path.join(os.path.dirname(__file__), "output", "Informatica_to_Fabric_Migration_Tool.pptx")
    prs.save(out_path)
    print(f"✅ Presentation saved: {out_path}")
    print(f"   {len(prs.slides)} slides generated")


if __name__ == "__main__":
    main()
